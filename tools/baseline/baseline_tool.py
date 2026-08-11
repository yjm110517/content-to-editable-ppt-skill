from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import time
import zipfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from PIL import Image


REPO_ROOT = Path(__file__).resolve().parents[2]
BASELINE_ROOT = REPO_ROOT / "baseline"
WORK_ROOT = REPO_ROOT / "work" / "p0-baseline"
RUNTIME_ROOT = REPO_ROOT / "content-to-editable-ppt"
RUNTIME_SCRIPTS = RUNTIME_ROOT / "scripts"
SCHEMA_ROOT = RUNTIME_ROOT / "schemas"
BASE_COMMIT = "dbd8f15fdfa7bc0c6558a68b6f4bf5cbd146d9e9"
CASES = ("B01", "B02", "B03", "B04", "B05", "B06")
TEXT_SUFFIXES = {".json", ".jsonl", ".log", ".md", ".txt", ".yaml", ".yml", ".xml", ".rels"}
SINGLE_FILE_BUDGET = 10 * 1024 * 1024
TOTAL_BUDGET = 100 * 1024 * 1024
PROJECT_AUTHOR = "content-to-editable-ppt-skill"


class BaselineError(RuntimeError):
    pass


def utc_now() -> str:
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat().replace("+00:00", "Z")


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def write_json(path: Path, value: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    temporary = path.with_name(f".{path.name}.tmp")
    temporary.write_text(json.dumps(value, ensure_ascii=False, indent=2, sort_keys=True) + "\n", encoding="utf-8")
    os.replace(temporary, path)


def load_json(path: Path) -> Any:
    return json.loads(path.read_text(encoding="utf-8"))


def command(argv: list[str], *, cwd: Path = REPO_ROOT, check: bool = True) -> subprocess.CompletedProcess[str]:
    completed = subprocess.run(argv, cwd=cwd, capture_output=True, text=True, encoding="utf-8", errors="replace", check=False)
    if check and completed.returncode != 0:
        raise BaselineError(f"command failed ({completed.returncode}): {' '.join(argv)}\n{completed.stdout}\n{completed.stderr}")
    return completed


def executable(candidate: str | None, fallback: str) -> Path:
    resolved = Path(candidate).resolve() if candidate else Path(shutil.which(fallback) or "").resolve()
    if not resolved.is_file():
        raise BaselineError(f"unable to resolve executable: {candidate or fallback}")
    return resolved


def version_tuple(text: str) -> tuple[int, ...]:
    match = re.search(r"(\d+)(?:\.(\d+))?(?:\.(\d+))?", text)
    if not match:
        raise BaselineError(f"unable to parse version from: {text!r}")
    return tuple(int(value or 0) for value in match.groups())


def runtime_tree_sha256() -> str:
    digest = hashlib.sha256()
    ignored = {"__pycache__", "node_modules", ".venv", ".agent-runtime"}
    for path in sorted(item for item in RUNTIME_ROOT.rglob("*") if item.is_file() and not ignored.intersection(item.parts)):
        relative = path.relative_to(RUNTIME_ROOT).as_posix()
        digest.update(relative.encode("utf-8"))
        digest.update(b"\0")
        digest.update(bytes.fromhex(sha256_file(path)))
    return digest.hexdigest()


def find_powerpoint() -> str:
    try:
        import winreg
    except ImportError as exc:
        raise BaselineError("PowerPoint detection requires Windows") from exc
    keys = [
        (winreg.HKEY_LOCAL_MACHINE, r"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths\POWERPNT.EXE"),
        (winreg.HKEY_LOCAL_MACHINE, r"SOFTWARE\WOW6432Node\Microsoft\Windows\CurrentVersion\App Paths\POWERPNT.EXE"),
        (winreg.HKEY_CURRENT_USER, r"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths\POWERPNT.EXE"),
    ]
    for hive, key in keys:
        try:
            with winreg.OpenKey(hive, key) as handle:
                value, _ = winreg.QueryValueEx(handle, None)
                if Path(value).is_file():
                    return value
        except OSError:
            continue
    raise BaselineError("Microsoft PowerPoint Desktop was not found")


def powerpoint_smoke(output_root: Path) -> dict[str, Any]:
    try:
        import win32com.client
    except ImportError as exc:
        raise BaselineError("pywin32 is required for PowerPoint COM smoke testing") from exc
    output_root.mkdir(parents=True, exist_ok=True)
    pptx = output_root / "powerpoint-com-smoke.pptx"
    png = output_root / "powerpoint-com-smoke.png"
    app = None
    presentation = None
    reopened = None
    started = time.perf_counter()
    try:
        app = win32com.client.DispatchEx("PowerPoint.Application")
        app.Visible = True
        presentation = app.Presentations.Add()
        slide = presentation.Slides.Add(1, 12)
        shape = slide.Shapes.AddTextbox(1, 72, 72, 420, 60)
        shape.TextFrame.TextRange.Text = "P0 PowerPoint COM Smoke Test"
        presentation.SaveAs(str(pptx), 24)
        slide.Export(str(png), "PNG", 640, 360)
        presentation.Close()
        presentation = None
        reopened = app.Presentations.Open(str(pptx), ReadOnly=True, Untitled=False, WithWindow=False)
        if reopened.Slides.Count != 1:
            raise BaselineError("PowerPoint COM smoke presentation has unexpected slide count")
        version = str(app.Version)
        reopened.Close()
        reopened = None
    finally:
        if reopened is not None:
            reopened.Close()
        if presentation is not None:
            presentation.Close()
        if app is not None:
            app.Quit()
    if not pptx.is_file() or not png.is_file():
        raise BaselineError("PowerPoint COM smoke did not produce PPTX and PNG outputs")
    return {
        "status": "pass",
        "powerpoint_version": version,
        "duration_ms": round((time.perf_counter() - started) * 1000),
        "pptx_sha256": sha256_file(pptx),
        "png_sha256": sha256_file(png),
    }


def environment(node_path: Path, python_path: Path) -> dict[str, Any]:
    node_version = command([str(node_path), "--version"]).stdout.strip()
    python_version = command([str(python_path), "--version"]).stdout.strip() or command([str(python_path), "--version"]).stderr.strip()
    if version_tuple(node_version) < (20, 0, 0):
        raise BaselineError(f"Node >=20 is required; got {node_version}")
    if version_tuple(python_version) < (3, 10, 0):
        raise BaselineError(f"Python >=3.10 is required; got {python_version}")
    powerpoint_path = find_powerpoint()
    smoke = powerpoint_smoke(WORK_ROOT / "_environment-smoke")
    pnpm = shutil.which("pnpm") or shutil.which("pnpm.cmd")
    if pnpm and Path(pnpm).suffix.lower() in {".cmd", ".bat"}:
        pnpm_version = command(["cmd.exe", "/c", pnpm, "--version"]).stdout.strip()
    else:
        pnpm_version = command([pnpm, "--version"]).stdout.strip() if pnpm else "unresolved"
    os_version = command(["cmd.exe", "/c", "ver"]).stdout.strip()
    windows_11_reference = "10.0.22000" in os_version or bool(re.search(r"10\.0\.(?:2[2-9]\d{3}|[3-9]\d{4})", os_version))
    return {
        "captured_at_utc": utc_now(),
        "source_runtime_commit": BASE_COMMIT,
        "runtime_tree_sha256": runtime_tree_sha256(),
        "operating_system": os_version,
        "reference_environment": {
            "target": "Windows 11 + Microsoft PowerPoint COM",
            "matches_target_os": windows_11_reference,
            "note": None if windows_11_reference else "Observed baseline host is Windows, but not the Windows 11 reference target.",
        },
        "node": {"version": node_version.lstrip("v"), "executable": "${NORMALIZED_RUNTIME_PATH}/node.exe"},
        "python": {"version": python_version.removeprefix("Python "), "executable": "${NORMALIZED_RUNTIME_PATH}/python.exe"},
        "pnpm": {"version": pnpm_version, "executable": "${NORMALIZED_RUNTIME_PATH}/pnpm" if pnpm else None},
        "powerpoint": {"executable": "${NORMALIZED_POWERPOINT_PATH}/POWERPNT.EXE", **smoke},
        "lockfiles": {
            "pnpm_lock_sha256": sha256_file(RUNTIME_SCRIPTS / "pnpm-lock.yaml"),
            "requirements_sha256": sha256_file(RUNTIME_SCRIPTS / "requirements.txt"),
        },
        "font_policy": ["Microsoft YaHei"],
        "iteration_policy": {"max_visual_revisions": 2, "max_total_iterations": 3},
    }


def prepare(case_id: str, node_path: Path, python_path: Path) -> dict[str, Any]:
    if case_id not in CASES:
        raise BaselineError(f"unknown case: {case_id}")
    case_root = BASELINE_ROOT / "cases" / case_id
    source = case_root / "input" / "source.png"
    request = case_root / "input" / "request.json"
    source_content = case_root / "evidence" / "baseline-source-content.json"
    for required in (source, request, source_content):
        if not required.is_file():
            raise BaselineError(f"missing baseline input: {required}")
    with Image.open(source) as image:
        if image.size != (1672, 941) or image.format != "PNG":
            raise BaselineError(f"{case_id} source must be a 1672x941 PNG")
    if load_json(source_content)["source_sha256"] != sha256_file(source):
        raise BaselineError(f"{case_id} source hash does not match baseline-source-content.json")
    case_work = WORK_ROOT / case_id
    if case_work.exists():
        raise BaselineError(f"work directory already exists; preserve or remove it explicitly before retrying: {case_work}")
    env_path = WORK_ROOT / "environment.json"
    if not env_path.exists():
        write_json(env_path, environment(node_path, python_path))
    case_work.mkdir(parents=True)
    shutil.copy2(source, case_work / "source.png")
    shutil.copy2(request, case_work / "request.json")
    shutil.copy2(source_content, case_work / "baseline-source-content.json")
    run_id = f"p0-{case_id.lower()}"
    log = case_work / "baseline.log"
    command([
        str(python_path), str(RUNTIME_SCRIPTS / "manage_run_state.py"), "init",
        "--work-root", str(case_work), "--request", str(case_work / "request.json"),
        "--output", str(case_work / "run_state.json"), "--run-id", run_id, "--log-file", str(log),
    ])
    command([
        str(python_path), str(RUNTIME_SCRIPTS / "manage_run_state.py"), "advance",
        "--work-root", str(case_work), "--state", str(case_work / "run_state.json"),
        "--event", "inputs_resolved", "--run-id", run_id, "--log-file", str(log),
    ])
    call_id = f"planner-{case_id.lower()}-i01"
    call_dir = case_work / ".agent-calls" / "01" / "planner" / call_id
    command([
        str(python_path), str(RUNTIME_SCRIPTS / "prepare_agent_call.py"),
        "--role", "planner", "--mode", "initial", "--work-root", str(case_work),
        "--request", str(case_work / "request.json"), "--source", str(case_work / "source.png"),
        "--iteration", "1", "--call-id", call_id, "--output-dir", str(call_dir),
        "--run-id", run_id, "--log-file", str(log),
    ])
    write_json(case_work / "baseline-session.json", {
        "case_id": case_id,
        "prepared_at_utc": utc_now(),
        "node_version": command([str(node_path), "--version"]).stdout.strip(),
        "python_version": command([str(python_path), "--version"]).stdout.strip(),
        "node_executable": "${NORMALIZED_RUNTIME_PATH}/node.exe",
        "python_executable": "${NORMALIZED_RUNTIME_PATH}/python.exe",
        "planner_call_dir": call_dir.relative_to(case_work).as_posix(),
        "max_visual_revisions": 2,
        "max_total_iterations": 3,
    })
    return {"status": "prepared", "case": case_id, "work_root": str(case_work), "planner_call_dir": str(call_dir)}


def normalize_text(text: str, case_work: Path) -> str:
    replacements = {
        str(REPO_ROOT): "${REPO_ROOT}",
        str(REPO_ROOT).replace("\\", "/"): "${REPO_ROOT}",
        str(case_work): "${CASE_ROOT}",
        str(case_work).replace("\\", "/"): "${CASE_ROOT}",
        str(Path.home()): "${USER_HOME}",
        str(Path.home()).replace("\\", "/"): "${USER_HOME}",
    }
    for source, target in sorted(replacements.items(), key=lambda item: len(item[0]), reverse=True):
        text = text.replace(source, target)
    return text


def copy_evidence(source: Path, target: Path, case_work: Path) -> None:
    target.parent.mkdir(parents=True, exist_ok=True)
    if source.suffix.lower() in TEXT_SUFFIXES:
        target.write_text(normalize_text(source.read_text(encoding="utf-8", errors="replace"), case_work), encoding="utf-8", newline="\n")
    else:
        shutil.copy2(source, target)


def local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


def sanitize_xml_metadata(name: str, data: bytes) -> bytes:
    if name not in {"docProps/core.xml", "docProps/app.xml"}:
        return data
    text = data.decode("utf-8")

    def replace_element(value: str, qualified_name: str, replacement: str) -> str:
        value = re.sub(
            rf"<{re.escape(qualified_name)}\b([^>]*)/>",
            lambda match: f"<{qualified_name}{match.group(1)}></{qualified_name}>",
            value,
            flags=re.IGNORECASE,
        )
        return re.sub(
            rf"(<{re.escape(qualified_name)}\b[^>]*>).*?(</{re.escape(qualified_name)}>)",
            lambda match: match.group(1) + replacement + match.group(2),
            value,
            flags=re.IGNORECASE | re.DOTALL,
        )

    if name == "docProps/core.xml":
        text = replace_element(text, "dc:creator", PROJECT_AUTHOR)
        text = replace_element(text, "cp:lastModifiedBy", PROJECT_AUTHOR)
    else:
        text = replace_element(text, "Company", "")
        text = replace_element(text, "Manager", "")
    return text.encode("utf-8")


def public_safety_findings(name: str, text: str) -> list[str]:
    patterns = {
        "windows_user_path": r"(?i)[A-Z]:\\Users\\",
        "local_drive_path": r"(?i)(?:^|[\"'>\s])[A-Z]:\\",
        "file_uri": r"(?i)file:///",
        "unc_path": r"\\\\[^\\\s]+\\",
        "email": r"(?i)\b[A-Z0-9._%+-]+@[A-Z0-9.-]+\.[A-Z]{2,}\b",
        "openai_key": r"\bsk-[A-Za-z0-9_-]{20,}\b",
        "github_token": r"\bgh[pousr]_[A-Za-z0-9]{20,}\b",
        "bearer_token": r"(?i)\bBearer\s+[A-Za-z0-9._-]{20,}",
    }
    findings = []
    for label, pattern in patterns.items():
        if re.search(pattern, text):
            findings.append(f"{name}: {label}")
    return findings


def sanitize_pptx(source: Path, target: Path, audit_path: Path) -> dict[str, Any]:
    target.parent.mkdir(parents=True, exist_ok=True)
    findings: list[str] = []
    external_relationships: list[str] = []
    with zipfile.ZipFile(source, "r") as archive, zipfile.ZipFile(target, "w") as output:
        for info in archive.infolist():
            data = sanitize_xml_metadata(info.filename, archive.read(info.filename))
            if info.filename.endswith((".xml", ".rels")):
                text = data.decode("utf-8", errors="replace")
                findings.extend(public_safety_findings(info.filename, text))
                if info.filename.endswith(".rels"):
                    import xml.etree.ElementTree as ET
                    try:
                        root = ET.fromstring(data)
                        for relation in root:
                            if relation.attrib.get("TargetMode") == "External":
                                external_relationships.append(f"{info.filename}: {relation.attrib.get('Target', '')}")
                    except ET.ParseError as exc:
                        findings.append(f"{info.filename}: invalid relationship XML: {exc}")
            output.writestr(info, data)
    with zipfile.ZipFile(target, "r") as check:
        corrupt = check.testzip()
        if corrupt:
            findings.append(f"corrupt ZIP entry: {corrupt}")
    if external_relationships:
        findings.extend(f"external_relationship: {item}" for item in external_relationships)
    audit = {
        "status": "pass" if not findings else "fail",
        "source_pptx_sha256": sha256_file(source),
        "sanitized_pptx_sha256": sha256_file(target),
        "creator": PROJECT_AUTHOR,
        "external_relationships": external_relationships,
        "findings": findings,
    }
    write_json(audit_path, audit)
    if findings:
        target.unlink(missing_ok=True)
        raise BaselineError("PPTX public-safety audit failed: " + "; ".join(findings))
    return audit


def render_sanitized_pptx(pptx: Path, output: Path) -> dict[str, Any]:
    try:
        import win32com.client
    except ImportError as exc:
        raise BaselineError("pywin32 is required to verify sanitized PPTX") from exc
    app = None
    presentation = None
    started = time.perf_counter()
    try:
        app = win32com.client.DispatchEx("PowerPoint.Application")
        app.Visible = True
        presentation = app.Presentations.Open(str(pptx), ReadOnly=True, Untitled=False, WithWindow=False)
        if presentation.Slides.Count != 1:
            raise BaselineError("sanitized baseline PPTX must contain exactly one slide")
        output.parent.mkdir(parents=True, exist_ok=True)
        presentation.Slides.Item(1).Export(str(output), "PNG", 1672, 941)
        presentation.Close()
        presentation = None
    finally:
        if presentation is not None:
            presentation.Close()
        if app is not None:
            app.Quit()
    return {"status": "pass", "duration_ms": round((time.perf_counter() - started) * 1000), "render_sha256": sha256_file(output)}


def extract_ppt_text(pptx: Path) -> list[str]:
    from pptx import Presentation
    presentation = Presentation(str(pptx))
    values: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()
                if text:
                    values.append(text)
    return values


def parse_log_timings(case_work: Path) -> dict[str, Any]:
    timestamps: list[datetime] = []
    components: dict[str, list[datetime]] = {}
    for log in case_work.rglob("*.log"):
        for line in log.read_text(encoding="utf-8", errors="ignore").splitlines():
            try:
                event = json.loads(line)
                raw = event.get("timestamp") or event.get("timestamp_utc") or event.get("at_utc")
                if not raw:
                    continue
                stamp = datetime.fromisoformat(raw.replace("Z", "+00:00"))
                timestamps.append(stamp)
                components.setdefault(str(event.get("component", "unknown")), []).append(stamp)
            except (json.JSONDecodeError, ValueError, TypeError):
                continue
    result: dict[str, Any] = {"total_observed_ms": None, "components_ms": {}}
    if timestamps:
        result["total_observed_ms"] = round((max(timestamps) - min(timestamps)).total_seconds() * 1000)
    for name, values in sorted(components.items()):
        result["components_ms"][name] = round((max(values) - min(values)).total_seconds() * 1000) if len(values) > 1 else 0
    return result


def count_technical_retries(case_work: Path) -> int:
    """Count repeated Runtime starts within one visual iteration.

    A failure before the Runtime emits ``started`` is a Host orchestration error,
    not a Runtime technical retry. Duplicate ``failed`` log events are therefore
    intentionally ignored.
    """
    retries = 0
    for iteration in sorted((case_work / "iterations").glob("[0-9][0-9]")):
        log = iteration / "pipeline.log"
        if not log.is_file():
            continue
        starts = 0
        for line in log.read_text(encoding="utf-8", errors="ignore").splitlines():
            try:
                event = json.loads(line)
            except json.JSONDecodeError:
                continue
            if event.get("component") == "run_pipeline" and event.get("event") == "started":
                starts += 1
        retries += max(0, starts - 1)
    return retries


def scan_public_safety(root: Path) -> list[str]:
    findings: list[str] = []
    for path in root.rglob("*"):
        if path.is_file() and path.suffix.lower() in TEXT_SUFFIXES:
            relative = path.relative_to(REPO_ROOT).as_posix()
            text = path.read_text(encoding="utf-8", errors="replace")
            findings.extend(public_safety_findings(relative, text))
    return findings


def capture(case_id: str) -> dict[str, Any]:
    if case_id not in CASES:
        raise BaselineError(f"unknown case: {case_id}")
    case_work = WORK_ROOT / case_id
    case_root = BASELINE_ROOT / "cases" / case_id
    if not case_work.is_dir():
        raise BaselineError(f"missing prepared work directory: {case_work}")
    iteration_root = case_work / "iterations"
    iterations = sorted(path for path in iteration_root.glob("[0-9][0-9]") if path.is_dir())
    if not iterations:
        raise BaselineError(f"{case_id} has no runtime iteration")
    numbers = [int(path.name) for path in iterations]
    if max(numbers) > 3 or (iteration_root / "04").exists():
        raise BaselineError(f"{case_id} violates max_total_iterations=3")
    final_source = iterations[-1]
    pptx_files = list(final_source.glob("*.pptx"))
    for required in ("layout.json", "crops.json", "asset_manifest.json", "qa_report.json", "rendered_slide.png"):
        if not (final_source / required).is_file():
            raise BaselineError(f"{case_id} final iteration is missing {required}")
    if len(pptx_files) != 1:
        raise BaselineError(f"{case_id} final iteration must contain exactly one PPTX")
    reviewer_success = (final_source / "review_report.json").is_file() and (final_source / "review_evaluation.json").is_file()
    reviewer_failure = (final_source / "reviewer_technical_failure.json").is_file()
    if not (reviewer_success or reviewer_failure):
        raise BaselineError(f"{case_id} is missing reviewer result or reviewer technical failure evidence")

    tracked_iterations = case_root / "evidence" / "iterations"
    if tracked_iterations.exists() and (case_root / "case-report.json").is_file():
        raise BaselineError(f"capture target already exists: {tracked_iterations}")
    lightweight = {
        "layout.json", "crops.json", "asset_manifest.json", "build_summary.json", "font_audit.json",
        "render_report.json", "qa_report.json", "review_report.json", "review_evaluation.json",
        "review_patch.json", "reviewer_technical_failure.json", "rendered_slide.png", "pipeline.log",
    }
    for iteration in iterations:
        target = tracked_iterations / iteration.name
        for name in lightweight:
            source = iteration / name
            if source.is_file():
                copy_evidence(source, target / name, case_work)
        call_root = case_work / ".agent-calls" / iteration.name
        if call_root.is_dir():
            for source in call_root.rglob("*"):
                if source.is_file() and source.name in {"call_manifest.json", "call_record.json"}:
                    relative = source.relative_to(call_root)
                    copy_evidence(source, target / "agent-calls" / relative, case_work)

    final_target = case_root / "evidence" / "final"
    audit_path = final_target / "pptx-public-safety-audit.json"
    sanitized_pptx = final_target / "baseline-final.pptx"
    audit = sanitize_pptx(pptx_files[0], sanitized_pptx, audit_path)
    render_check = render_sanitized_pptx(sanitized_pptx, final_target / "sanitized-render-check.png")
    assets_source = final_source / "assets"
    if assets_source.is_dir():
        for source in assets_source.rglob("*"):
            if source.is_file():
                copy_evidence(source, final_target / "assets" / source.relative_to(assets_source), case_work)
    copy_evidence(case_work / "run_state.json", case_root / "evidence" / "run_state.json", case_work)
    copy_evidence(case_work / "baseline-session.json", case_root / "evidence" / "baseline-session.json", case_work)
    prior_attempt_path = case_work / "baseline-prior-attempt.json"
    prior_issues: list[dict[str, Any]] = []
    if prior_attempt_path.is_file():
        copy_evidence(prior_attempt_path, case_root / "evidence" / "baseline-prior-attempt.json", case_work)
        prior_issues = load_json(prior_attempt_path).get("known_issues", [])

    expected = [item["text"] for item in load_json(case_root / "evidence" / "baseline-source-content.json")["text_items"]]
    actual = extract_ppt_text(sanitized_pptx)
    combined = "\n".join(actual)
    missing = [text for text in expected if text not in combined]
    call_records = list(tracked_iterations.rglob("call_record.json"))
    planner_calls = sum(1 for path in call_records if "planner" in path.parts)
    reviewer_calls = sum(1 for path in call_records if "reviewer" in path.parts)
    manifest = load_json(final_source / "asset_manifest.json")
    crops = load_json(final_source / "crops.json")
    asset_count = len(manifest.get("assets", []))
    crop_count = len(crops.get("assets", []))
    runtime_asset_count = sum(1 for path in assets_source.rglob("*") if path.is_file()) if assets_source.is_dir() else 0
    qa = load_json(final_source / "qa_report.json")
    review_state = "technical_failure" if reviewer_failure else load_json(final_source / "review_evaluation.json").get("policy_decision", "unknown")
    technical_retries = count_technical_retries(case_work)
    report = {
        "schema_version": "p0-v1",
        "case_id": case_id,
        "captured_at_utc": utc_now(),
        "source_runtime_commit": BASE_COMMIT,
        "final_iteration": int(final_source.name),
        "max_visual_revisions": 2,
        "max_total_iterations": 3,
        "counts": {
            "planner_calls": planner_calls,
            "reviewer_calls": reviewer_calls,
            "technical_retries": technical_retries,
            "visual_revisions": max(0, int(final_source.name) - 1),
            "asset_count": asset_count,
            "crop_entry_count": crop_count,
            "runtime_asset_count": runtime_asset_count,
        },
        "outcome": {"structural_qa": qa.get("status"), "visual_review": review_state},
        "content_comparison": {"expected_count": len(expected), "ppt_text_shape_count": len(actual), "missing_text": missing},
        "timings": parse_log_timings(case_work),
        "public_safety": {**audit, "sanitized_render_verification": render_check},
        "artifacts": {
            "source_sha256": sha256_file(case_root / "input" / "source.png"),
            "final_pptx_sha256": sha256_file(sanitized_pptx),
            "final_render_sha256": sha256_file(final_target / "sanitized-render-check.png"),
        },
        "known_issues": ([{"category": "content", "description": f"Missing frozen text: {text}"} for text in missing]
                         + ([{"category": "reviewer_technical_failure", "description": "Visual Reviewer did not return a valid report; structured failure evidence was captured."}] if reviewer_failure else [])
                         + ([{"category": "runtime_technical_retry", "description": "Runtime required a same-iteration technical retry after an initial pipeline start failed."}] if technical_retries else [])
                         + prior_issues),
    }
    write_json(case_root / "case-report.json", report)
    findings = scan_public_safety(case_root)
    if findings:
        raise BaselineError("public-safety scan failed: " + "; ".join(findings))
    rebuild_manifest_and_report()
    enforce_budgets()
    return {"status": "captured", "case": case_id, "final_iteration": int(final_source.name), "reviewer": review_state}


def rebuild_manifest_and_report() -> None:
    reports = []
    for case_id in CASES:
        path = BASELINE_ROOT / "cases" / case_id / "case-report.json"
        if path.is_file():
            reports.append(load_json(path))
    environment_path = WORK_ROOT / "environment.json"
    manifest = {
        "baseline_version": "p0-v1",
        "source_runtime_commit": BASE_COMMIT,
        "runtime_tree_sha256": runtime_tree_sha256(),
        "environment": load_json(environment_path) if environment_path.is_file() else None,
        "iteration_policy": {"max_visual_revisions": 2, "max_total_iterations": 3},
        "cases": [{
            "case_id": item["case_id"], "final_iteration": item["final_iteration"],
            "outcome": item["outcome"], "counts": item["counts"], "artifacts": item["artifacts"],
        } for item in reports],
    }
    write_json(BASELINE_ROOT / "manifest.json", manifest)
    lines = [
        "# P0 Baseline Freeze Report", "", f"- Baseline: `p0-v1`", f"- Runtime commit: `{BASE_COMMIT}`",
        "- Authority: [Architecture](../docs/architecture/v2.0/overall-architecture-and-development-plan.md), "
        "[Test Plan](../docs/testing/v1.0/test-and-acceptance-plan.md), [ADR](../DECISIONS.md)", "",
        "| Case | Final iteration | Structural QA | Visual review | Planner | Reviewer | Technical retries | Revisions | Known issues |", "|---|---:|---|---|---:|---:|---:|---:|---:|",
    ]
    for item in reports:
        counts = item["counts"]
        lines.append(
            f"| [{item['case_id']}](cases/{item['case_id']}/case-report.json) | {item['final_iteration']} | "
            f"{item['outcome']['structural_qa']} | {item['outcome']['visual_review']} | {counts['planner_calls']} | "
            f"{counts['reviewer_calls']} | {counts['technical_retries']} | {counts['visual_revisions']} | {len(item['known_issues'])} |"
        )
    lines.extend(["", "This report freezes observed behavior. It is not a pixel-level golden and does not claim that current visual issues are fixed.", ""])
    (BASELINE_ROOT / "baseline-report.md").write_text("\n".join(lines), encoding="utf-8", newline="\n")


def enforce_budgets() -> None:
    total = 0
    oversized = []
    for path in BASELINE_ROOT.rglob("*"):
        if path.is_file():
            total += path.stat().st_size
            if path.stat().st_size > SINGLE_FILE_BUDGET:
                oversized.append(f"{path.relative_to(BASELINE_ROOT)} ({path.stat().st_size} bytes)")
    if oversized:
        raise BaselineError("single-file baseline budget exceeded: " + ", ".join(oversized))
    if total > TOTAL_BUDGET:
        raise BaselineError(f"aggregate baseline budget exceeded: {total} > {TOTAL_BUDGET}")


def verify_markdown_links() -> list[str]:
    failures: list[str] = []
    files = [REPO_ROOT / "README.md"] + list(BASELINE_ROOT.rglob("*.md"))
    pattern = re.compile(r"\[[^\]]+\]\(([^)]+)\)")
    for path in files:
        for target in pattern.findall(path.read_text(encoding="utf-8")):
            clean = target.split("#", 1)[0]
            if not clean:
                continue
            if re.match(r"(?i)^(?:file://|/[A-Z]:|[A-Z]:[\\/])", clean):
                failures.append(f"{path.relative_to(REPO_ROOT)}: absolute local link {target}")
            elif not re.match(r"^[a-z]+://", clean) and not (path.parent / clean).resolve().exists():
                failures.append(f"{path.relative_to(REPO_ROOT)}: missing link {target}")
    return failures


def verify() -> dict[str, Any]:
    failures: list[str] = []
    manifest_path = BASELINE_ROOT / "manifest.json"
    if not manifest_path.is_file():
        failures.append("baseline/manifest.json is missing")
        manifest = {"cases": []}
    else:
        manifest = load_json(manifest_path)
    if manifest.get("source_runtime_commit") != BASE_COMMIT:
        failures.append("manifest source Runtime commit is incorrect")
    if {item.get("case_id") for item in manifest.get("cases", [])} != set(CASES):
        failures.append("manifest must contain exactly B01-B06")
    diff = command(["git", "diff", "--exit-code", BASE_COMMIT, "--", "content-to-editable-ppt"], check=False)
    if diff.returncode != 0:
        failures.append("Runtime directory changed relative to dbd8f15")
    for case_id in CASES:
        case_root = BASELINE_ROOT / "cases" / case_id
        required = [
            case_root / "input" / "source.png", case_root / "input" / "request.json",
            case_root / "evidence" / "baseline-source-content.json", case_root / "case-report.json",
        ]
        for path in required:
            if not path.is_file():
                failures.append(f"{case_id}: missing {path.relative_to(case_root)}")
        if not all(path.is_file() for path in required):
            continue
        report = load_json(case_root / "case-report.json")
        final_iteration = int(report.get("final_iteration", 0))
        if final_iteration not in {1, 2, 3} or (case_root / "evidence" / "iterations" / "04").exists():
            failures.append(f"{case_id}: invalid iteration count")
        final_iteration_root = case_root / "evidence" / "iterations" / f"{final_iteration:02d}"
        final_root = case_root / "evidence" / "final"
        for path in (final_iteration_root / "layout.json", final_iteration_root / "qa_report.json", final_iteration_root / "rendered_slide.png", final_root / "baseline-final.pptx", final_root / "pptx-public-safety-audit.json"):
            if not path.is_file():
                failures.append(f"{case_id}: missing final evidence {path.name}")
        reviewer_ok = (final_iteration_root / "review_report.json").is_file() and (final_iteration_root / "review_evaluation.json").is_file()
        reviewer_failed = (final_iteration_root / "reviewer_technical_failure.json").is_file()
        if not (reviewer_ok or reviewer_failed):
            failures.append(f"{case_id}: Reviewer invocation evidence missing")
        if report.get("counts", {}).get("reviewer_calls", 0) < 1 and not reviewer_failed:
            failures.append(f"{case_id}: Reviewer call count is zero")
        manifest_file = final_iteration_root / "asset_manifest.json"
        crops_file = final_iteration_root / "crops.json"
        if manifest_file.is_file() and crops_file.is_file():
            asset_count = len(load_json(manifest_file).get("assets", []))
            crop_count = len(load_json(crops_file).get("assets", []))
            runtime_assets = report.get("counts", {}).get("runtime_asset_count")
            if case_id in {"B02", "B03", "B05"} and (asset_count < 1 or runtime_assets < 1):
                failures.append(f"{case_id}: expected non-empty Runtime Assets")
            if case_id == "B06" and (asset_count != 0 or crop_count != 0 or runtime_assets != 0):
                failures.append("B06: zero-asset semantics violated")
        audit = final_root / "pptx-public-safety-audit.json"
        if audit.is_file() and load_json(audit).get("status") != "pass":
            failures.append(f"{case_id}: PPTX public-safety audit did not pass")
    failures.extend(verify_markdown_links())
    failures.extend(scan_public_safety(BASELINE_ROOT))
    try:
        enforce_budgets()
    except BaselineError as exc:
        failures.append(str(exc))
    if failures:
        raise BaselineError("P0 Baseline verification failed:\n- " + "\n- ".join(failures))
    return {"status": "pass", "cases": list(CASES), "baseline_bytes": sum(path.stat().st_size for path in BASELINE_ROOT.rglob("*") if path.is_file())}


def main() -> int:
    parser = argparse.ArgumentParser(description="Prepare, capture, and verify P0 Baseline evidence.")
    subparsers = parser.add_subparsers(dest="command", required=True)
    prepare_parser = subparsers.add_parser("prepare")
    prepare_parser.add_argument("--case", required=True, choices=CASES)
    prepare_parser.add_argument("--node-path")
    prepare_parser.add_argument("--python-path")
    capture_parser = subparsers.add_parser("capture")
    capture_parser.add_argument("--case", required=True, choices=CASES)
    subparsers.add_parser("verify")
    args = parser.parse_args()
    try:
        if args.command == "prepare":
            node_path = executable(args.node_path, "node")
            python_path = executable(args.python_path, "python")
            result = prepare(args.case, node_path, python_path)
        elif args.command == "capture":
            result = capture(args.case)
        else:
            result = verify()
        print(json.dumps(result, ensure_ascii=False, indent=2))
        return 0
    except Exception as exc:
        print(json.dumps({"status": "error", "error": str(exc)}, ensure_ascii=False, indent=2), file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
