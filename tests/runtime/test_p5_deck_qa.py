from __future__ import annotations

import json
import tempfile
import unittest
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

ROOT = Path(__file__).resolve().parents[2]
sys_path = str(ROOT / "content-to-editable-ppt" / "scripts")
import sys
if sys_path not in sys.path:
    sys.path.insert(0, sys_path)

from deck_qa import run_deck_qa  # noqa: E402
from schema_utils import ContractError  # noqa: E402


def _make_pptx(path: Path, slides: int = 2, *, full_picture: bool = False) -> None:
    """Minimal deterministic PPTX built with python-pptx."""
    from pptx import Presentation
    from pptx.util import Emu
    presentation = Presentation()
    for index in range(slides):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        if full_picture:
            from pptx.shapes.picture import Picture
            picture = slide.shapes.add_picture
            stream = tempfile.TemporaryFile(suffix=".png")
            from PIL import Image
            image = Image.new("RGB", (32, 32), "white")
            image.save(stream, format="PNG")
            stream.seek(0)
            from pptx.util import Inches
            slide.shapes.add_picture(stream, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        else:
            box = slide.shapes.add_textbox(Emu(100000), Emu(100000), Emu(1000000), Emu(200000))
            box.text_frame.text = f"Slide {index + 1} text"
    presentation.save(str(path))


def _add_external_relationship(path: Path) -> None:
    """Inject an external relationship into the first slide rels."""
    with zipfile.ZipFile(path, "r") as archive:
        names = archive.namelist()
        rel_name = next(name for name in names if name.startswith("ppt/slides/_rels/slide1"))
        rels = archive.read(rel_name)
        tree = ET.fromstring(rels)
        ns = {"pr": "http://schemas.openxmlformats.org/package/2006/relationships"}
        element = ET.SubElement(tree, f"{{{ns['pr']}}}Relationship")
        element.set("Id", "rIdExt1")
        element.set("Type", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink")
        element.set("Target", "https://example.com/external")
        element.set("TargetMode", "External")
        new_rels = ET.tostring(tree, encoding="unicode")
        content = {}
        for name in names:
            content[name] = archive.read(name)
        content[rel_name] = new_rels.encode("utf-8")
    with zipfile.ZipFile(path, "w") as archive:
        for name, data in sorted(content.items()):
            archive.writestr(name, data)


def _p4_bundle(deck_id: str, candidate_sha: str, manifest_sha: str) -> dict:
    return {
        "candidate_pptx_sha256": candidate_sha,
        "reconstruction_manifest_sha256": manifest_sha,
        "status": "pass",
        "post_assembly_slide_drift": 0,
        "unexpected_assembly_mutation": 0,
    }


class P5DeckQATests(unittest.TestCase):
    def test_clean_deck_passes(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            pptx = root / "candidate.pptx"
            _make_pptx(pptx, slides=2)
            import hashlib
            candidate_sha = hashlib.sha256(pptx.read_bytes()).hexdigest()
            manifest_sha = "a" * 64
            report = run_deck_qa(
                deck_id="D05",
                candidate_pptx=pptx,
                p4_state={"state": "p4_complete", "deck_id": "D05", "current_artifacts": {"candidate_deck_sha256": candidate_sha, "reconstruction_manifest_sha256": manifest_sha}},
                p4_manifest={"deck_id": "D05", "slides": [{"order": 1, "slide_id": "S01"}, {"order": 2, "slide_id": "S02"}]},
                p4_candidate_report={"candidate_pptx_sha256": candidate_sha, "reconstruction_manifest_sha256": manifest_sha},
                p4_drift_report={"candidate_deck_sha256": candidate_sha, "status": "pass", "post_assembly_slide_drift": 0, "unexpected_assembly_mutation": 0},
            )
            self.assertEqual(report["status"], "pass")
            self.assertEqual(report["blocking_issues"], 0)
            self.assertEqual(report["exception_pages"], [])
            self.assertEqual(report["content_drift"], 0)

    def test_hash_chain_break_is_critical(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            pptx = root / "candidate.pptx"
            _make_pptx(pptx, slides=1)
            report = run_deck_qa(
                deck_id="D05",
                candidate_pptx=pptx,
                p4_state={"state": "p4_complete", "deck_id": "D05", "current_artifacts": {"candidate_deck_sha256": "0" * 64, "reconstruction_manifest_sha256": "1" * 64}},
                p4_manifest={"deck_id": "D05", "slides": [{"order": 1, "slide_id": "S01"}]},
                p4_candidate_report={"candidate_pptx_sha256": "2" * 64, "reconstruction_manifest_sha256": "1" * 64},
                p4_drift_report={"candidate_deck_sha256": "0" * 64, "status": "pass", "post_assembly_slide_drift": 0, "unexpected_assembly_mutation": 0},
            )
            self.assertEqual(report["status"], "fail")
            self.assertTrue(any(item["code"] == "p4_hash_chain" and item["severity"] == "critical" for item in report["issues"]))

    def test_full_slide_raster_substitution_detected(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            pptx = root / "candidate.pptx"
            _make_pptx(pptx, slides=1, full_picture=True)
            import hashlib
            candidate_sha = hashlib.sha256(pptx.read_bytes()).hexdigest()
            manifest_sha = "a" * 64
            report = run_deck_qa(
                deck_id="D05",
                candidate_pptx=pptx,
                p4_state={"state": "p4_complete", "deck_id": "D05", "current_artifacts": {"candidate_deck_sha256": candidate_sha, "reconstruction_manifest_sha256": manifest_sha}},
                p4_manifest={"deck_id": "D05", "slides": [{"order": 1, "slide_id": "S01"}]},
                p4_candidate_report={"candidate_pptx_sha256": candidate_sha, "reconstruction_manifest_sha256": manifest_sha},
                p4_drift_report={"candidate_deck_sha256": candidate_sha, "status": "pass", "post_assembly_slide_drift": 0, "unexpected_assembly_mutation": 0},
            )
            self.assertTrue(any(item["code"] == "full_slide_raster_substitution" for item in report["issues"]))
            self.assertIn("S01", report["exception_pages"])

    def test_external_relationship_is_blocking(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            pptx = root / "candidate.pptx"
            _make_pptx(pptx, slides=1)
            _add_external_relationship(pptx)
            import hashlib
            candidate_sha = hashlib.sha256(pptx.read_bytes()).hexdigest()
            manifest_sha = "a" * 64
            report = run_deck_qa(
                deck_id="D05",
                candidate_pptx=pptx,
                p4_state={"state": "p4_complete", "deck_id": "D05", "current_artifacts": {"candidate_deck_sha256": candidate_sha, "reconstruction_manifest_sha256": manifest_sha}},
                p4_manifest={"deck_id": "D05", "slides": [{"order": 1, "slide_id": "S01"}]},
                p4_candidate_report={"candidate_pptx_sha256": candidate_sha, "reconstruction_manifest_sha256": manifest_sha},
                p4_drift_report={"candidate_deck_sha256": candidate_sha, "status": "pass", "post_assembly_slide_drift": 0, "unexpected_assembly_mutation": 0},
            )
            self.assertTrue(any(item["code"] == "external_relationship" for item in report["issues"]))
            self.assertEqual(report["status"], "fail")
            self.assertGreater(report["blocking_issues"], 0)

    def test_internal_relationship_missing_target_is_blocking(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            pptx = root / "candidate.pptx"
            _make_pptx(pptx, slides=1)
            with zipfile.ZipFile(pptx, "r") as archive:
                names = archive.namelist()
                rel_name = next(name for name in names if name.startswith("ppt/slides/_rels/slide1"))
                content = {}
                for name in names:
                    content[name] = archive.read(name)
                rels = content[rel_name].decode("utf-8")
                rels = rels.replace("</Relationships>", '<Relationship Id="rIdBroken" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="media/missing.png" TargetMode="Internal"/></Relationships>')
                content[rel_name] = rels.encode("utf-8")
            with zipfile.ZipFile(pptx, "w") as archive:
                for name, data in sorted(content.items()):
                    archive.writestr(name, data)
            import hashlib
            candidate_sha = hashlib.sha256(pptx.read_bytes()).hexdigest()
            manifest_sha = "a" * 64
            report = run_deck_qa(
                deck_id="D05",
                candidate_pptx=pptx,
                p4_state={"state": "p4_complete", "deck_id": "D05", "current_artifacts": {"candidate_deck_sha256": candidate_sha, "reconstruction_manifest_sha256": manifest_sha}},
                p4_manifest={"deck_id": "D05", "slides": [{"order": 1, "slide_id": "S01"}]},
                p4_candidate_report={"candidate_pptx_sha256": candidate_sha, "reconstruction_manifest_sha256": manifest_sha},
                p4_drift_report={"candidate_deck_sha256": candidate_sha, "status": "pass", "post_assembly_slide_drift": 0, "unexpected_assembly_mutation": 0},
            )
            self.assertTrue(any(item["code"] == "unsafe_path_relationship" for item in report["issues"]))
            self.assertEqual(report["status"], "fail")

    def test_macro_content_type_is_blocking(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            pptx = root / "candidate.pptx"
            _make_pptx(pptx, slides=1)
            with zipfile.ZipFile(pptx, "r") as archive:
                names = archive.namelist()
                content = {}
                for name in names:
                    content[name] = archive.read(name)
                ct = content["[Content_Types].xml"].decode("utf-8")
                ct = ct.replace("</Types>", '<Override PartName="/ppt/vbaProject.bin" ContentType="application/vnd.ms-office.vbaProject"/></Types>')
                content["[Content_Types].xml"] = ct.encode("utf-8")
                content["ppt/vbaProject.bin"] = b"vba"
            with zipfile.ZipFile(pptx, "w") as archive:
                for name, data in sorted(content.items()):
                    archive.writestr(name, data)
            import hashlib
            candidate_sha = hashlib.sha256(pptx.read_bytes()).hexdigest()
            manifest_sha = "a" * 64
            report = run_deck_qa(
                deck_id="D05",
                candidate_pptx=pptx,
                p4_state={"state": "p4_complete", "deck_id": "D05", "current_artifacts": {"candidate_deck_sha256": candidate_sha, "reconstruction_manifest_sha256": manifest_sha}},
                p4_manifest={"deck_id": "D05", "slides": [{"order": 1, "slide_id": "S01"}]},
                p4_candidate_report={"candidate_pptx_sha256": candidate_sha, "reconstruction_manifest_sha256": manifest_sha},
                p4_drift_report={"candidate_deck_sha256": candidate_sha, "status": "pass", "post_assembly_slide_drift": 0, "unexpected_assembly_mutation": 0},
            )
            self.assertTrue(any(item["code"] == "macro_ole" for item in report["issues"]))
            self.assertEqual(report["status"], "fail")


if __name__ == "__main__":
    unittest.main()
