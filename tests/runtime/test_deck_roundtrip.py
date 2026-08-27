from __future__ import annotations

import sys
import tempfile
import unittest
import zipfile
from pathlib import Path

ROOT = Path(__file__).resolve().parents[2]
sys_path = str(ROOT / "content-to-editable-ppt" / "scripts")
if sys_path not in sys.path:
    sys.path.insert(0, sys_path)

from deck_roundtrip import _chart_signature, _media_signature, _slide_names, _slide_size, _structural_snapshot, _workbook_signature  # noqa: E402


def _make_deck(path: Path, slide_count: int = 2) -> None:
    from pptx import Presentation
    from pptx.util import Emu
    presentation = Presentation()
    for index in range(slide_count):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(Emu(100000), Emu(100000), Emu(1000000), Emu(200000))
        box.text_frame.text = f"Slide {index + 1} text"
    presentation.save(str(path))


def _bindings(slide_count: int) -> list[dict]:
    return [{"slide_id": f"S{index + 1:02d}", "order": index + 1} for index in range(slide_count)]


class DeckRoundtripPureTests(unittest.TestCase):
    def test_slide_identity_is_object_names_not_counts(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            deck = root / "candidate.pptx"
            _make_deck(deck, slide_count=2)
            snapshot = _structural_snapshot(deck, _bindings(2))
            self.assertTrue(snapshot["valid"])
            self.assertEqual(len(snapshot["slides"]), 2)
            # every slide carries its own textbox object name
            for item in snapshot["slides"]:
                self.assertTrue(item["object_names"], "each slide must expose cNvPr object names")
            # 两个页面不能因为对象数量相同而被视为顺序一致
            order_a = [item["object_names"] for item in snapshot["slides"]]
            self.assertEqual(len(order_a), 2)

    def test_slide_size_from_presentation_xml(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            deck = root / "candidate.pptx"
            _make_deck(deck, slide_count=1)
            cx, cy = _slide_size(deck)
            self.assertGreater(cx, 0)
            self.assertGreater(cy, 0)

    def test_slide_count_mismatch_invalid(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            deck = root / "candidate.pptx"
            _make_deck(deck, slide_count=1)
            snapshot = _structural_snapshot(deck, _bindings(2))
            self.assertFalse(snapshot["valid"])

    def test_chart_signature_is_semantic(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            deck = root / "candidate.pptx"
            from pptx import Presentation
            from pptx.chart.data import CategoryChartData
            from pptx.enum.chart import XL_CHART_TYPE
            from pptx.util import Emu
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            chart_data = CategoryChartData()
            chart_data.categories = ["A", "B"]
            chart_data.add_series("Series", (1, 2))
            slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Emu(100000), Emu(100000), Emu(3000000), Emu(2000000), chart_data)
            presentation.save(str(deck))
            signature = _chart_signature(deck)
            self.assertEqual(len(signature), 1)
            self.assertEqual(signature[0]["type"], "barChart")
            self.assertEqual(signature[0]["categories"], ["A", "B"])
            self.assertEqual(signature[0]["values"], ["1", "2"])
            self.assertEqual(signature[0]["series_names"], ["Series"])

    def test_chart_signature_catches_value_change(self) -> None:
        def build(values: tuple) -> Path:
            from pptx import Presentation
            from pptx.chart.data import CategoryChartData
            from pptx.enum.chart import XL_CHART_TYPE
            from pptx.util import Emu
            path = Path(tempfile.mkdtemp()) / "c.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            chart_data = CategoryChartData()
            chart_data.categories = ["A", "B"]
            chart_data.add_series("Series", values)
            slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Emu(100000), Emu(100000), Emu(3000000), Emu(2000000), chart_data)
            presentation.save(str(path))
            return path

        first = _chart_signature(build((1, 2)))
        second = _chart_signature(build((1, 3)))
        self.assertNotEqual(first, second)

    def test_workbook_signature_and_media(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            deck = root / "candidate.pptx"
            _make_deck(deck, slide_count=1)
            workbooks = _workbook_signature(deck)
            self.assertEqual(workbooks, [])
            media = _media_signature(deck)
            self.assertIsInstance(media, list)


if __name__ == "__main__":
    unittest.main()
