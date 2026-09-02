from __future__ import annotations

import json
import os
import subprocess
import sys
import tempfile
import unittest
import zipfile
from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[2]
SCRIPTS = ROOT / "content-to-editable-ppt" / "scripts"
sys.path.insert(0, str(SCRIPTS))

from tests.runtime.test_direct_deck_request import request
from slide_size import resolve_slide_size


class DirectDeckBuilderTests(unittest.TestCase):
    def test_builder_creates_normalized_multi_page_native_deck(self):
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp); value = request()
            second = json.loads(json.dumps(value["slides"][0])); second["slide_id"] = "S02"; second["order"] = 2
            second["elements"] = [
                {"id": "TITLE", "type": "text", "x": 0.6, "y": 0.4, "w": 5.0, "h": 0.7, "z_index": 10, "editable": True, "text": "Chart", "content_ref": "S02-TITLE", "style_ref": "body"},
                {"id": "CHART", "type": "chart", "x": 6.0, "y": 1.4, "w": 5.5, "h": 4.5, "z_index": 5, "editable": True, "chart_type": "vertical_bar", "categories": ["A", "B"], "series": [{"name": "Values", "values": [1, 2]}], "value_scale": 1},
            ]
            value["slides"].append(second)
            request_path = root / "request.json"; request_path.write_text(json.dumps(value), encoding="utf-8")
            assets = root / "assets.json"; assets.write_text("[]", encoding="utf-8")
            output, report = root / "deck.pptx", root / "build-report.json"
            command = [os.environ.get("DECK_TEST_NODE", "node"), str(SCRIPTS / "build_deck.mjs"), "--request", str(request_path), "--staged-assets", str(assets), "--output", str(output), "--report", str(report)]
            completed = subprocess.run(command, cwd=ROOT, capture_output=True, text=True, encoding="utf-8", errors="replace")
            self.assertEqual(completed.returncode, 0, completed.stdout + completed.stderr)
            built = json.loads(report.read_text(encoding="utf-8")); self.assertEqual(built["slide_order"], ["S01", "S02"])
            presentation = Presentation(output); self.assertEqual(len(presentation.slides), 2)
            expected_size = resolve_slide_size(value["output_ratio"])
            self.assertAlmostEqual(expected_size["width_in"], presentation.slide_width / 914400, places=3)
            self.assertAlmostEqual(expected_size["height_in"], presentation.slide_height / 914400, places=3)
            names = [[shape.name for shape in slide.shapes] for slide in presentation.slides]
            self.assertIn("ivt:TITLE", names[0]); self.assertIn("ivt:TITLE", names[1]); self.assertIn("ivt:CHART", names[1])
            with zipfile.ZipFile(output) as archive:
                self.assertTrue(any(name.startswith("ppt/charts/chart") for name in archive.namelist()))


if __name__ == "__main__":
    unittest.main()
